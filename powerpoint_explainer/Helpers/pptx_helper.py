from pptx import Presentation


def parse_presentation(pptx_path):
    """
    Parses the .pptx file and returns a list of slides.

    Args:
        pptx_path (str): Path to the .pptx file.

    Returns:
        list: List of slides.
    """
    presentation = Presentation(pptx_path)
    slides = presentation.slides
    return slides


def extract_text(slide):
    """
    Extracts the text from a slide's text boxes and returns a list of text strings.

    Args:
        slide (Slide): PowerPoint slide object.

    Returns:
        list: List of extracted text strings.
    """
    text_list = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_list.append(text)
    return text_list
